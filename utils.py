# utils.py
import psycopg2
import time
from typing import List
from PyPDF2 import PdfReader
from langchain.docstore.document import Document
import bcrypt
import base64
import docx
import pytesseract
from PIL import Image
import os
from pdf2image import convert_from_path
import logging

# Set up logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[logging.FileHandler('app.log'), logging.StreamHandler()]
)
logger = logging.getLogger(__name__)

HAS_PDF2IMAGE = True  # Assuming pdf2image is installed as per the try-except block intent
try:
    from pptx import Presentation
except ImportError:
    logging.warning("python-pptx not found. PPTX support will be limited.")

def process_pdf(uploaded_file) -> List[Document]:
    """Process a PDF file and return a list of Document objects."""
    pdf_reader = PdfReader(uploaded_file)
    documents = []
    for page_num, page in enumerate(pdf_reader.pages):
        text = page.extract_text()
        if not text.strip() and HAS_PDF2IMAGE and os.path.exists(uploaded_file.name):
            try:
                images = convert_from_path(uploaded_file.name, first_page=page_num + 1, last_page=page_num + 1)
                text = pytesseract.image_to_string(images[0]) if images else ""
            except Exception as e:
                logger.warning(f"OCR failed for page {page_num + 1}: {e}")
                text = ""
        if text:
            documents.append(Document(
                page_content=text,
                metadata={"page": page_num + 1, "filename": uploaded_file.name}
            ))
    return documents

def process_docx(uploaded_file) -> List[Document]:
    """Process a DOCX file and return a list of Document objects."""
    doc = docx.Document(uploaded_file)
    text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
    return [Document(page_content=text, metadata={"filename": uploaded_file.name})]

def process_pptx(uploaded_file) -> List[Document]:
    """Process a PPTX file and return a list of Document objects."""
    try:
        ppt = Presentation(uploaded_file)
        full_text = "\n".join([shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text_frame") and shape.text_frame])
        return [Document(page_content=full_text, metadata={"filename": uploaded_file.name})]
    except Exception as e:
        logger.warning(f"PPTX processing failed: {e}")
        return []

def process_txt(uploaded_file) -> List[Document]:
    """Process a TXT file and return a list of Document objects."""
    text = uploaded_file.read().decode("utf-8")
    return [Document(page_content=text, metadata={"filename": uploaded_file.name})]

def process_image(uploaded_file) -> List[Document]:
    """Process an image file and return a list of Document objects using OCR."""
    try:
        image = Image.open(uploaded_file)
        text = pytesseract.image_to_string(image)
        if text:
            return [Document(page_content=text, metadata={"filename": uploaded_file.name})]
        return []
    except Exception as e:
        logger.warning(f"Image OCR failed: {e}")
        return []

def process_attachment(uploaded_file):
    """Process various file types and return a list of Document objects."""
    file_type = uploaded_file.type.lower()
    filename = uploaded_file.name.lower()
    
    if file_type == "application/pdf" or filename.endswith('.pdf'):
        return process_pdf(uploaded_file)
    elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document" or filename.endswith('.docx'):
        return process_docx(uploaded_file)
    elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation" or filename.endswith('.pptx'):
        return process_pptx(uploaded_file)
    elif file_type == "text/plain" or filename.endswith('.txt'):
        return process_txt(uploaded_file)
    elif file_type in ["image/jpeg", "image/png"] or filename.endswith(('.jpg', '.jpeg', '.png')):
        return process_image(uploaded_file)
    else:
        logger.warning(f"Unsupported file type: {file_type} ({filename})")
        return []

def get_db_connection():
    """Get the database connection from the session state (defined in app.py)."""
    from streamlit.runtime.scriptrunner import get_script_run_ctx
    ctx = get_script_run_ctx()
    if ctx and hasattr(ctx.session_state, 'db_connection'):
        return ctx.session_state.db_connection
    raise Exception("Database connection not found in session state")

def init_database():
    """Initialize the database with proper schema and handle migration."""
    conn = get_db_connection()
    if conn is None:
        logger.error("Failed to establish database connection")
        raise Exception("Database connection is not available")
    c = conn.cursor()

    # Create or update users table
    c.execute('''CREATE TABLE IF NOT EXISTS users (
                 id SERIAL PRIMARY KEY,
                 username TEXT UNIQUE NOT NULL,
                 password TEXT NOT NULL,
                 display_name TEXT)''')

    # Check and migrate chat_history table
    c.execute("SELECT column_name FROM information_schema.columns WHERE table_name = 'chat_history'")
    columns = [col[0] for col in c.fetchall()]
    c.execute("SELECT table_name FROM information_schema.tables WHERE table_schema = 'public'")
    tables = c.fetchall()
    if "chat_history" in [table[0] for table in tables]:
        if "file_sources" not in columns:
            c.execute("ALTER TABLE chat_history RENAME TO chat_history_old")
            c.execute('''CREATE TABLE chat_history (
                id SERIAL PRIMARY KEY,
                username TEXT NOT NULL,
                chat_id INTEGER NOT NULL,
                timestamp TEXT NOT NULL,
                user_message TEXT NOT NULL,
                bot_response TEXT NOT NULL,
                file_sources TEXT)''')
            c.execute("INSERT INTO chat_history (username, chat_id, timestamp, user_message, bot_response) SELECT username, chat_id, timestamp, user_message, bot_response FROM chat_history_old")
            c.execute("DROP TABLE chat_history_old")
            conn.commit()
            logger.info("Migrated chat_history table to include file_sources column")
    else:
        c.execute('''CREATE TABLE IF NOT EXISTS chat_history (
                     id SERIAL PRIMARY KEY,
                     username TEXT NOT NULL,
                     chat_id INTEGER NOT NULL,
                     timestamp TEXT NOT NULL,
                     user_message TEXT NOT NULL,
                     bot_response TEXT NOT NULL,
                     file_sources TEXT)''')

    # Check and migrate user_activity table
    c.execute("SELECT column_name FROM information_schema.columns WHERE table_name = 'user_activity'")
    columns = [col[0] for col in c.fetchall()]
    c.execute("SELECT table_name FROM information_schema.tables WHERE table_schema = 'public'")
    tables = c.fetchall()
    if "user_activity" in [table[0] for table in tables]:
        if "details" not in columns:
            c.execute("ALTER TABLE user_activity RENAME TO user_activity_old")
            c.execute('''CREATE TABLE user_activity (
                id SERIAL PRIMARY KEY,
                username TEXT NOT NULL,
                activity_type TEXT NOT NULL,
                details TEXT,
                timestamp TEXT DEFAULT CURRENT_TIMESTAMP)''')
            c.execute("INSERT INTO user_activity (username, activity_type, timestamp) SELECT username, activity_type, timestamp FROM user_activity_old")
            c.execute("DROP TABLE user_activity_old")
            conn.commit()
            logger.info("Migrated user_activity table to include details column")
    else:
        c.execute('''CREATE TABLE IF NOT EXISTS user_activity (
                     id SERIAL PRIMARY KEY,
                     username TEXT NOT NULL,
                     activity_type TEXT NOT NULL,
                     details TEXT,
                     timestamp TEXT DEFAULT CURRENT_TIMESTAMP)''')

    # Check and migrate file_processing table
    c.execute("SELECT column_name FROM information_schema.columns WHERE table_name = 'file_processing'")
    columns = [col[0] for col in c.fetchall()]
    c.execute("SELECT table_name FROM information_schema.tables WHERE table_schema = 'public'")
    tables = c.fetchall()
    if "file_processing" in [table[0] for table in tables]:
        if "size" not in columns:
            c.execute("ALTER TABLE file_processing RENAME TO file_processing_old")
            c.execute('''CREATE TABLE file_processing (
                id SERIAL PRIMARY KEY,
                username TEXT NOT NULL,
                filename TEXT NOT NULL,
                size INTEGER,
                status TEXT,
                timestamp TEXT DEFAULT CURRENT_TIMESTAMP)''')
            c.execute("INSERT INTO file_processing (username, filename, status, timestamp) SELECT username, filename, status, timestamp FROM file_processing_old")
            c.execute("DROP TABLE file_processing_old")
            conn.commit()
            logger.info("Migrated file_processing table to include size column")
    else:
        c.execute('''CREATE TABLE IF NOT EXISTS file_processing (
                     id SERIAL PRIMARY KEY,
                     username TEXT NOT NULL,
                     filename TEXT NOT NULL,
                     size INTEGER,
                     status TEXT,
                     timestamp TEXT DEFAULT CURRENT_TIMESTAMP)''')

    conn.commit()

def login_user_base64(username: str, password: str) -> bool:
    """Authenticate a user with base64-encoded hashed password."""
    conn = get_db_connection()
    c = conn.cursor()
    c.execute("SELECT password FROM users WHERE username = %s", (username,))
    result = c.fetchone()
    if result:
        hashed_b64 = result[0]
        try:
            hashed_password = base64.b64decode(hashed_b64.encode('utf-8'))
            return bcrypt.checkpw(password.encode('utf-8'), hashed_password)
        except Exception as e:
            logger.error(f"Password verification error: {e}")
            return False
    return False

def register_user_base64(username: str, password: str, display_name: str = None) -> bool:
    """Register a new user with base64-encoded hashed password and display name."""
    conn = get_db_connection()
    c = conn.cursor()
    hashed = bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt())
    hashed_b64 = base64.b64encode(hashed).decode('utf-8')
    try:
        c.execute("INSERT INTO users (username, password, display_name) VALUES (%s, %s, %s)", (username, hashed_b64, display_name))
        conn.commit()
        return True
    except psycopg2.IntegrityError:
        return False

def save_chat_history(username: str, user_message: str, bot_response: str, chat_id: int = 1, file_sources: list = None):
    """Save chat history to the database with chat_id and file_sources."""
    conn = get_db_connection()
    c = conn.cursor()
    timestamp = time.strftime("%Y-%m-%d %H:%M:%S")
    c.execute("INSERT INTO chat_history (username, chat_id, timestamp, user_message, bot_response, file_sources) VALUES (%s, %s, %s, %s, %s, %s)",
              (username, chat_id, timestamp, user_message, bot_response, str(file_sources) if file_sources else None))
    conn.commit()

def get_chat_history(username: str, chat_id: int) -> List[dict]:
    """Retrieve chat history for a user and specific chat_id."""
    conn = get_db_connection()
    c = conn.cursor()
    c.execute("SELECT user_message, bot_response, timestamp, file_sources FROM chat_history WHERE username = %s AND chat_id = %s ORDER BY timestamp",
              (username, chat_id))
    history = [{"user_message": row[0], "bot_response": row[1], "timestamp": row[2], "file_sources": eval(row[3]) if row[3] else []} for row in c.fetchall()]
    return history

def get_user_chats(username: str) -> List[int]:
    """Retrieve distinct chat IDs for a user."""
    conn = get_db_connection()
    c = conn.cursor()
    c.execute("SELECT DISTINCT chat_id FROM chat_history WHERE username = %s ORDER BY chat_id DESC", (username,))
    chat_ids = [row[0] for row in c.fetchall()]
    return chat_ids

def delete_chat_history(username: str, chat_id: int):
    """Delete chat history for a specific chat_id."""
    conn = get_db_connection()
    c = conn.cursor()
    c.execute("DELETE FROM chat_history WHERE username = %s AND chat_id = %s", (username, chat_id))
    conn.commit()
    logger.info(f"Deleted chat history for chat_id: {chat_id} for user: {username}")

def log_user_activity(username: str, activity_type: str, details: str = None):
    """Log user activity."""
    conn = get_db_connection()
    c = conn.cursor()
    c.execute("INSERT INTO user_activity (username, activity_type, details) VALUES (%s, %s, %s)",
              (username, activity_type, details))
    conn.commit()

def log_file_processing(username: str, filename: str, size: int, status: str):
    """Log file processing details."""
    conn = get_db_connection()
    c = conn.cursor()
    c.execute("INSERT INTO file_processing (username, filename, size, status) VALUES (%s, %s, %s, %s)",
              (username, filename, size, status))
    conn.commit()